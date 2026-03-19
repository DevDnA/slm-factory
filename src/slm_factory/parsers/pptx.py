"""PPTX (PowerPoint) 문서 파서입니다."""

from __future__ import annotations

from pathlib import Path
from typing import ClassVar

from ..models import ParsedDocument
from ..utils import get_logger
from .base import BaseParser, extract_date_from_filename, rows_to_markdown

logger = get_logger("parsers.pptx")


class PPTXParser(BaseParser):
    """PPTX 프레젠테이션을 파싱합니다."""

    extensions: ClassVar[list[str]] = [".pptx"]

    def can_parse_content(self, path: Path) -> bool:
        import zipfile

        try:
            if not zipfile.is_zipfile(str(path)):
                return False
            with zipfile.ZipFile(path, "r") as zf:
                return "ppt/presentation.xml" in zf.namelist()
        except Exception:
            return False

    def parse(self, path: Path) -> ParsedDocument:
        """PPTX 파일을 파싱하여 ParsedDocument를 반환합니다.

        슬라이드 텍스트, 표 (마크다운), 발표자 노트를 추출합니다.

        매개변수
        ----------
        path:
            PPTX 파일의 경로입니다.

        반환값
        -------
        ParsedDocument
            텍스트 내용, 표 (마크다운), 메타데이터로 채워집니다.

        예외
        ------
        ImportError
            python-pptx가 설치되지 않았으면 발생합니다.
        FileNotFoundError
            *path*가 존재하지 않으면 발생합니다.
        RuntimeError
            PPTX 파일을 파싱할 수 없으면 발생합니다.
        """
        try:
            from pptx import Presentation
            from pptx.table import Table as PptxTable
        except ImportError as exc:
            raise ImportError(
                "python-pptx가 설치되지 않았습니다. uv sync --extra pptx 로 설치하세요."
            ) from exc

        path = Path(path)
        if not path.exists():
            raise FileNotFoundError(f"PPTX file not found: {path}")

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise RuntimeError(
                f"PPTX 파일을 파싱할 수 없습니다: {path}\n"
                f"원인: 파일이 손상되었거나 python-pptx가 지원하지 않는 형식일 수 있습니다.\n"
                f"해결: PowerPoint에서 다시 저장하거나 PDF로 변환하세요."
            ) from exc

        # ------------------------------------------------------------------
        # 슬라이드 텍스트 및 표 추출
        # ------------------------------------------------------------------
        slide_texts: list[str] = []
        tables_md: list[str] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            parts: list[str] = [f"## 슬라이드 {slide_idx}"]

            for shape in slide.shapes:
                # 표 처리
                if shape.has_table:
                    table: PptxTable = shape.table
                    rows: list[list[str]] = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        md_table = rows_to_markdown(rows)
                        if md_table:
                            tables_md.append(md_table)
                            parts.append(md_table)
                    continue

                # 텍스트 프레임 처리
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            parts.append(text)

            # 발표자 노트 추출
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f"\n> **노트**: {notes_text}")

            if len(parts) > 1:  # 제목 외에 내용이 있을 때만
                slide_texts.append("\n\n".join(parts))

        content = "\n\n".join(slide_texts)

        # ------------------------------------------------------------------
        # 제목 / 메타데이터
        # ------------------------------------------------------------------
        title = path.stem
        if prs.core_properties and prs.core_properties.title:
            title = prs.core_properties.title

        metadata: dict[str, object] = {
            "slide_count": len(prs.slides),
        }
        if prs.core_properties:
            if prs.core_properties.author:
                metadata["author"] = prs.core_properties.author
            if prs.core_properties.created:
                metadata["date"] = prs.core_properties.created.strftime("%Y-%m-%d")

        date_from_name = extract_date_from_filename(path.stem)
        if date_from_name and "date" not in metadata:
            metadata["date"] = date_from_name

        return ParsedDocument(
            doc_id=path.name,
            title=title,
            content=content,
            tables=tables_md,
            metadata=metadata,
        )
